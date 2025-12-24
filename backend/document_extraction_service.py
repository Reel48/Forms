"""
Document Extraction Service
Enhanced text extraction for knowledge base documents
Supports PDF, Excel, PowerPoint, and DOCX files
"""
import io
import logging
from typing import Optional
import mimetypes

logger = logging.getLogger(__name__)

DEFAULT_MAX_TEXT_CHARS = 100000  # 100KB for knowledge base documents (larger than chat attachments)


def extract_text_from_excel(data: bytes, max_chars: int = DEFAULT_MAX_TEXT_CHARS) -> str:
    """
    Extract text from Excel files (.xlsx, .xls)
    
    Args:
        data: Excel file bytes
        max_chars: Maximum characters to extract
        
    Returns:
        Extracted text with sheet names and cell values
    """
    try:
        try:
            import openpyxl
            from openpyxl import load_workbook
        except ImportError:
            logger.warning("openpyxl not installed - Excel extraction unavailable")
            return ""
        
        workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        parts = []
        
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            sheet_parts = [f"Sheet: {sheet_name}"]
            
            # Extract cell values
            for row in sheet.iter_rows(values_only=True):
                row_values = []
                for cell_value in row:
                    if cell_value is not None:
                        # Convert to string, handle different types
                        if isinstance(cell_value, (int, float)):
                            row_values.append(str(cell_value))
                        else:
                            row_values.append(str(cell_value).strip())
                
                if row_values:
                    sheet_parts.append(" | ".join(row_values))
            
            if len(sheet_parts) > 1:  # More than just the sheet name
                parts.append("\n".join(sheet_parts))
        
        text = "\n\n".join(parts).strip()
        return text[:max_chars]
        
    except Exception as e:
        logger.warning(f"Excel text extraction failed: {str(e)}")
        return ""


def extract_text_from_powerpoint(data: bytes, max_chars: int = DEFAULT_MAX_TEXT_CHARS) -> str:
    """
    Extract text from PowerPoint files (.pptx, .ppt)
    
    Args:
        data: PowerPoint file bytes
        max_chars: Maximum characters to extract
        
    Returns:
        Extracted text with slide numbers, titles, and content
    """
    try:
        try:
            from pptx import Presentation
        except ImportError:
            logger.warning("python-pptx not installed - PowerPoint extraction unavailable")
            return ""
        
        presentation = Presentation(io.BytesIO(data))
        parts = []
        
        for slide_num, slide in enumerate(presentation.slides, 1):
            slide_parts = [f"Slide {slide_num}:"]
            
            # Extract title (if exists)
            if slide.shapes.title:
                title_text = slide.shapes.title.text.strip()
                if title_text:
                    slide_parts.append(f"Title: {title_text}")
            
            # Extract text from all shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text = shape.text.strip()
                    if text and text not in slide_parts:  # Avoid duplicates
                        slide_parts.append(text)
            
            # Extract notes (if available)
            if hasattr(slide, "notes_slide") and slide.notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    slide_parts.append(f"Notes: {notes_text}")
            
            if len(slide_parts) > 1:  # More than just "Slide N:"
                parts.append("\n".join(slide_parts))
        
        text = "\n\n".join(parts).strip()
        return text[:max_chars]
        
    except Exception as e:
        logger.warning(f"PowerPoint text extraction failed: {str(e)}")
        return ""


def extract_text_from_pdf_enhanced(data: bytes, max_chars: int = DEFAULT_MAX_TEXT_CHARS) -> str:
    """
    Enhanced PDF extraction with better structure preservation
    
    Args:
        data: PDF file bytes
        max_chars: Maximum characters to extract
        
    Returns:
        Extracted text with page markers
    """
    try:
        try:
            from pypdf import PdfReader
        except ImportError:
            try:
                from PyPDF2 import PdfReader
            except ImportError:
                logger.warning("pypdf or PyPDF2 not installed - PDF extraction unavailable")
                return ""
        
        reader = PdfReader(io.BytesIO(data))
        parts = []
        
        for page_num, page in enumerate(reader.pages[:100], 1):  # Limit to 100 pages
            try:
                page_text = page.extract_text()
                if page_text and page_text.strip():
                    # Add page marker for context
                    parts.append(f"Page {page_num}:\n{page_text.strip()}")
            except Exception:
                continue
        
        text = "\n\n".join(parts).strip()
        return text[:max_chars]
        
    except Exception as e:
        logger.warning(f"PDF text extraction failed: {str(e)}")
        return ""


def normalize_text(text: str) -> str:
    """
    Clean and normalize extracted text
    
    Args:
        text: Raw extracted text
        
    Returns:
        Normalized text
    """
    if not text:
        return ""
    
    # Remove excessive whitespace
    import re
    text = re.sub(r'\n{3,}', '\n\n', text)  # Max 2 consecutive newlines
    text = re.sub(r'[ \t]+', ' ', text)  # Multiple spaces to single space
    
    # Remove control characters except newlines and tabs
    text = ''.join(char for char in text if ord(char) >= 32 or char in '\n\t')
    
    return text.strip()


def extract_text_from_document(
    data: bytes,
    mime_type: Optional[str] = None,
    file_name: Optional[str] = None,
    max_chars: int = DEFAULT_MAX_TEXT_CHARS
) -> str:
    """
    Main function to extract text from any supported document type
    
    Args:
        data: File bytes
        mime_type: MIME type of the file
        file_name: Name of the file
        max_chars: Maximum characters to extract
        
    Returns:
        Extracted and normalized text
    """
    mime = (mime_type or "").lower()
    name = (file_name or "").lower()
    
    # Guess MIME type if not provided
    if not mime and name:
        mime, _ = mimetypes.guess_type(name)
        mime = (mime or "").lower()
    
    # Plain text / CSV
    if mime.startswith("text/") or name.endswith((".txt", ".csv")):
        try:
            text = data.decode("utf-8", errors="replace")
            return normalize_text(text[:max_chars])
        except Exception:
            return ""
    
    # PDF
    if "pdf" in mime or name.endswith(".pdf"):
        text = extract_text_from_pdf_enhanced(data, max_chars)
        return normalize_text(text)
    
    # Excel
    if ("excel" in mime or "spreadsheet" in mime or 
        name.endswith((".xlsx", ".xls"))):
        text = extract_text_from_excel(data, max_chars)
        return normalize_text(text)
    
    # PowerPoint
    if ("presentation" in mime or "powerpoint" in mime or
        name.endswith((".pptx", ".ppt"))):
        text = extract_text_from_powerpoint(data, max_chars)
        return normalize_text(text)
    
    # DOCX (Word)
    if "word" in mime or name.endswith(".docx"):
        try:
            import docx  # type: ignore
            doc = docx.Document(io.BytesIO(data))
            text = "\n".join([p.text for p in doc.paragraphs if p.text]).strip()
            return normalize_text(text[:max_chars])
        except Exception as e:
            logger.warning(f"DOCX text extraction failed: {str(e)}")
            return ""
    
    logger.warning(f"Unsupported file type: {mime} / {name}")
    return ""

